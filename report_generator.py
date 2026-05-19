"""Generates a professional 10-slide PowerPoint report from competitor data."""
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# Brand colours
DARK  = RGBColor(0x0D,0x1B,0x2A)
BLUE  = RGBColor(0x1A,0x73,0xE8)
ACCENT= RGBColor(0x00,0xC9,0xA7)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LIGHT = RGBColor(0xF0,0xF4,0xFF)
GREY  = RGBColor(0x6B,0x72,0x80)

W = Inches(13.33)
H = Inches(7.5)

def _solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color

def _txt(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tf.text = text
    p = tf.paragraphs[0]; p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text; run.font.size = Pt(size)
    run.font.bold = bold; run.font.color.rgb = color
    run.font.name = "Calibri"

def _add_textbox(slide, text, l, t, w, h, size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    _txt(tf, str(text), size, bold, color, align)
    return tb

def _bg(slide, color=DARK):
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    _solid(bg, color); bg.line.fill.background(); bg.zorder = 0

def _header_bar(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(1.2))
    _solid(bar, BLUE); bar.line.fill.background()
    _add_textbox(slide, title, 0.3, 0.1, 10, 0.7, 28, True)
    if subtitle:
        _add_textbox(slide, subtitle, 0.3, 0.75, 10, 0.4, 13, False, RGBColor(0xCC,0xDD,0xFF))

def _rank_companies(all_data):
    """Score companies: subs(30) + avg_views(30) + freq(20) + engagement(20)"""
    scores = {}
    def norm(vals):
        mx = max(vals) if max(vals) else 1
        return [v/mx for v in vals]
    names = [d["company_name"] for d in all_data]
    subs  = norm([d["subscriber_count"] for d in all_data])
    views = norm([d["avg_views"] for d in all_data])
    freq  = norm([d["upload_freq_per_week"] for d in all_data])
    eng   = norm([sum(v["engagement_rate"] for v in d["videos"])/max(len(d["videos"]),1) for d in all_data])
    for i,n in enumerate(names):
        scores[n] = round(subs[i]*30 + views[i]*30 + freq[i]*20 + eng[i]*20, 1)
    return scores

def _fmt(n):
    if n>=1_000_000: return f"{n/1_000_000:.1f}M"
    if n>=1_000:     return f"{n/1_000:.1f}K"
    return str(int(n))

# ── Slides ──────────────────────────────────────────────────────────────────

def slide_cover(prs, all_data, your_company):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, DARK)
    accent = sl.shapes.add_shape(1, 0, Inches(2.8), Inches(0.5), Inches(1.9))
    _solid(accent, ACCENT); accent.line.fill.background()
    _add_textbox(sl,"VIDEO COMPETITOR INTELLIGENCE REPORT",0.7,1.0,12,1.1,32,True,WHITE,PP_ALIGN.LEFT)
    names = " · ".join(d["company_name"] for d in all_data)
    _add_textbox(sl,f"Companies analysed: {names}",0.7,2.9,12,0.5,14,False,RGBColor(0xAA,0xCC,0xFF))
    _add_textbox(sl,f"Your company: {your_company}",0.7,3.5,12,0.5,13,False,ACCENT)
    _add_textbox(sl,f"Report date: {datetime.now().strftime('%d %B %Y')}",0.7,4.1,12,0.4,12,False,GREY)
    _add_textbox(sl,"Powered by MyPromoVideos Intelligence Engine",0.7,6.8,12,0.4,10,False,GREY)

def slide_executive_summary(prs, all_data, scores):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl); _header_bar(sl,"Executive Summary","Who is leading in video marketing and why")
    ranked = sorted(scores.items(), key=lambda x:x[1], reverse=True)
    leader = ranked[0][0]
    leader_data = next((d for d in all_data if d["company_name"]==leader), all_data[0])
    bullets = [
        f"🏆 Video Marketing Leader: {leader} (Score: {ranked[0][1]}/100)",
        f"   Subscribers: {_fmt(leader_data['subscriber_count'])}  |  Avg Views: {_fmt(leader_data['avg_views'])}  |  Upload Freq: {leader_data['upload_freq_per_week']:.2f}/week",
        "",
    ]
    for name, score in ranked[1:]:
        d = next((x for x in all_data if x["company_name"]==name), None)
        if d:
            bullets.append(f"• {name} — Score {score}/100 | Subs: {_fmt(d['subscriber_count'])} | Avg Views: {_fmt(d['avg_views'])}")
    bullets += ["","Key Insight: The leader dominates through consistent uploads, high subscriber engagement,","and a focused content strategy. Gaps exist in short-form and educational content."]
    y = 1.4
    for b in bullets:
        _add_textbox(sl, b, 0.5, y, 12.3, 0.35, 12 if b.startswith(" ") else 13, b.startswith("🏆"), WHITE if not b.startswith(" ") else RGBColor(0xAA,0xCC,0xFF))
        y += 0.33

def slide_channel_overview(prs, all_data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl); _header_bar(sl,"Channel Overview","Subscriber count · Total videos · Upload frequency")
    headers = ["Company","Channel","Subscribers","Total Videos","Total Views","Upload/Week","Country"]
    col_w   = [2.2, 2.5, 1.5, 1.5, 1.6, 1.4, 1.1]
    x_pos   = [0.3]; 
    for w in col_w[:-1]: x_pos.append(x_pos[-1]+w)
    # header row
    for i,(h,x,w) in enumerate(zip(headers,x_pos,col_w)):
        _add_textbox(sl,h,x,1.35,w-0.1,0.35,11,True,ACCENT)
    # data rows
    for ri, d in enumerate(all_data):
        y = 1.8 + ri*0.65
        row_bg = sl.shapes.add_shape(1,Inches(0.3),Inches(y-0.05),Inches(12.73),Inches(0.6))
        _solid(row_bg, RGBColor(0x1A,0x2A,0x3A) if ri%2==0 else RGBColor(0x14,0x22,0x30))
        row_bg.line.fill.background()
        vals = [d["company_name"], d.get("channel_title","—"), _fmt(d["subscriber_count"]),
                _fmt(d["video_count"]), _fmt(d["view_count"]),
                str(d["upload_freq_per_week"]), d.get("country","N/A")]
        for val,x,w in zip(vals,x_pos,col_w):
            _add_textbox(sl,val,x,y,w-0.1,0.55,10,False,WHITE)

def slide_content_performance(prs, all_data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl); _header_bar(sl,"Content Performance","Top performing videos by views and engagement")
    y = 1.4
    for d in all_data:
        _add_textbox(sl, d["company_name"], 0.3, y, 12, 0.3, 13, True, ACCENT)
        y += 0.35
        if d["top_videos"]:
            for v in d["top_videos"][:3]:
                line = f"▶  {v['title'][:70]}{'…' if len(v['title'])>70 else ''}   |   👁 {_fmt(v['views'])}   👍 {_fmt(v['likes'])}   💬 {_fmt(v['comments'])}   ⚡ {v['engagement_rate']}%"
                _add_textbox(sl, line, 0.5, y, 12.3, 0.32, 10, False, WHITE)
                y += 0.33
        else:
            _add_textbox(sl,"No video data available.",0.5,y,12,0.3,10,False,GREY)
            y += 0.33
        y += 0.1

def slide_topics_themes(prs, all_data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl); _header_bar(sl,"Content Topics & Themes","What each company covers and what they are missing")
    col_w = 13.3/len(all_data)
    for ci, d in enumerate(all_data):
        x = 0.15 + ci*col_w
        box = sl.shapes.add_shape(1,Inches(x),Inches(1.3),Inches(col_w-0.2),Inches(5.8))
        _solid(box, RGBColor(0x1A,0x2A,0x3A)); box.line.fill.background()
        _add_textbox(sl, d["company_name"], x+0.1, 1.4, col_w-0.3, 0.35, 12, True, ACCENT)
        _add_textbox(sl,"Top Topics:", x+0.1, 1.85, col_w-0.3, 0.3, 10, True, WHITE)
        y = 2.2
        for t in d.get("topics",[])[:8]:
            _add_textbox(sl, f"• {t}", x+0.15, y, col_w-0.4, 0.3, 10, False, WHITE)
            y += 0.32
        if not d.get("topics"):
            _add_textbox(sl,"No data",x+0.1,y,col_w-0.3,0.3,10,False,GREY)

def slide_posting_frequency(prs, all_data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl); _header_bar(sl,"Posting Frequency & Consistency","Who is most active and on what cadence")
    mx = max((d["upload_freq_per_week"] for d in all_data),default=1) or 1
    
    # Elegant sizing with plenty of breathing room at the top (avoiding header collision)
    bar_h = 3.6         # Max height of bars (reduced from 4.5 to fit labels nicely)
    bar_area_top = 2.2  # Top boundary (pushed down from 1.5 to leave 1.0 inch below header bar)
    bar_w = 1.2         # Width of each bar
    
    total_w = len(all_data) * (bar_w + 0.4)
    start_x = (13.33 - total_w) / 2
    
    # Draw a thin elegant baseline for the bar chart
    baseline = sl.shapes.add_shape(1, Inches(start_x - 0.2), Inches(bar_area_top + bar_h), Inches(total_w), Inches(0.02))
    _solid(baseline, GREY)
    baseline.line.fill.background()
    
    for i, d in enumerate(all_data):
        x = start_x + i * (bar_w + 0.4)
        val = d["upload_freq_per_week"]
        h = max(0.2, (val / mx) * bar_h)
        top = bar_area_top + (bar_h - h)
        
        # Add the bar shape
        bar = sl.shapes.add_shape(1, Inches(x), Inches(top), Inches(bar_w), Inches(h))
        colors = [BLUE, ACCENT, RGBColor(0xFF,0x6B,0x6B), RGBColor(0xFF,0xD9,0x3D), RGBColor(0xA2,0x9B,0xF5)]
        _solid(bar, colors[i % len(colors)])
        bar.line.fill.background()
        
        # Format label (use decimal only if it is small, e.g. < 10)
        val_lbl = f"{val:.2f}/wk" if val < 10 else f"{int(val)}/wk"
        
        # Value label above the bar (increased vertical spacing and uses Pt size 12)
        _add_textbox(sl, val_lbl, x, top - 0.42, bar_w, 0.35, 12, True, WHITE, PP_ALIGN.CENTER)
        
        # Company name below the bar (with generous vertical padding)
        _add_textbox(sl, d["company_name"], x, bar_area_top + bar_h + 0.12, bar_w, 0.4, 11, True, WHITE, PP_ALIGN.CENTER)

def slide_engagement(prs, all_data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl); _header_bar(sl,"Engagement Analysis","Average views · likes · comments per video")
    metrics=[("Avg Views","avg_views"),("Avg Likes","avg_likes"),("Avg Comments","avg_comments")]
    col_w=13.0/len(all_data); 
    for ci,d in enumerate(all_data):
        x=0.15+ci*col_w
        hd=sl.shapes.add_shape(1,Inches(x),Inches(1.3),Inches(col_w-0.15),Inches(0.45))
        _solid(hd,BLUE); hd.line.fill.background()
        _add_textbox(sl,d["company_name"],x+0.1,1.35,col_w-0.3,0.35,12,True,WHITE)
        y=1.9
        for label,key in metrics:
            card=sl.shapes.add_shape(1,Inches(x),Inches(y),Inches(col_w-0.15),Inches(1.2))
            _solid(card,RGBColor(0x1A,0x2A,0x3A)); card.line.fill.background()
            _add_textbox(sl,label,x+0.1,y+0.1,col_w-0.3,0.3,10,False,GREY)
            _add_textbox(sl,_fmt(d[key]),x+0.1,y+0.4,col_w-0.3,0.6,24,True,ACCENT)
            y+=1.35

def slide_gap_analysis(prs, all_data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl); _header_bar(sl,"Gap Analysis","Topics and formats competitors are not covering")
    all_topics=set()
    for d in all_data: all_topics.update(d.get("topics",[]))
    gaps=[]
    potential=["case study","tutorial","behind the scenes","testimonial","product demo",
               "faq","industry news","comparison","tips","how-to","live stream","podcast"]
    for p in potential:
        covered=sum(1 for d in all_data if any(p in t for t in d.get("topics",[])))
        if covered<len(all_data)//2+1:
            gaps.append(p)
    _add_textbox(sl,"Identified Content Gaps Across All Competitors:",0.4,1.35,12,0.35,13,True,ACCENT)
    y=1.8
    for g in gaps[:8]:
        _add_textbox(sl,f"◆  '{g.title()}' — underutilised across {len(all_data)-sum(1 for d in all_data if any(g in t for t in d.get('topics',[])))}/{len(all_data)} competitors",0.5,y,12.3,0.32,12,False,WHITE)
        y+=0.37
    _add_textbox(sl,"→ Opportunity: First-mover advantage in these formats can capture uncontested audience segments.",0.4,y+0.2,12.3,0.4,12,True,ACCENT)

def slide_recommendations(prs, all_data, your_company):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl); _header_bar(sl,"Video Marketing Recommendations",f"Actionable steps for {your_company}")
    recs=[
        ("1. Increase Upload Cadence","Aim for at least 2 videos/week. Consistency drives algorithmic reach and subscriber retention."),
        ("2. Double Down on High-Engagement Formats","Tutorials and how-to videos consistently outperform promotional content. Prioritise educational value."),
        ("3. Fill Content Gaps","Produce case studies and behind-the-scenes content — competitors are largely absent from these formats."),
        ("4. Optimise for Search (YouTube SEO)","Use keyword-rich titles, timestamps in descriptions, and relevant tags to improve discoverability."),
        ("5. Leverage Shorts for Reach","YouTube Shorts offer outsized reach with minimal production cost. Repurpose top-performing clips."),
        ("6. Engage With Comments","Reply to comments within 24 hours. Channels with active comment engagement see 40% higher retention."),
        ("7. Collaborate & Cross-Promote","Partner with micro-influencers in your niche to tap new audiences at low cost."),
    ]
    y=1.35
    for title,body in recs:
        _add_textbox(sl,title,0.4,y,12.5,0.3,12,True,ACCENT)
        y+=0.32
        _add_textbox(sl,body,0.6,y,12.2,0.3,11,False,WHITE)
        y+=0.38

def slide_summary_ranking(prs, all_data, scores):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl); _header_bar(sl,"Summary & Rankings","Overall scorecard across all companies")
    ranked=sorted(scores.items(),key=lambda x:x[1],reverse=True)
    medals=["🥇","🥈","🥉","4️⃣","5️⃣"]
    cols=["Rank","Company","Score /100","Subscribers","Avg Views","Upload/Week","Engagement"]
    col_w=[0.8,2.5,1.4,1.6,1.5,1.5,1.5]
    x_pos=[0.2]
    for w in col_w[:-1]: x_pos.append(x_pos[-1]+w)
    for col,x,w in zip(cols,x_pos,col_w):
        _add_textbox(sl,col,x,1.35,w-0.1,0.35,11,True,ACCENT,PP_ALIGN.CENTER)
    for ri,(name,score) in enumerate(ranked):
        d=next((x for x in all_data if x["company_name"]==name),None)
        if not d: continue
        y=1.82+ri*0.72
        rb=sl.shapes.add_shape(1,Inches(0.2),Inches(y-0.07),Inches(13.0),Inches(0.65))
        _solid(rb,RGBColor(0x1A,0x2A,0x3A) if ri%2==0 else RGBColor(0x14,0x22,0x30))
        rb.line.fill.background()
        avg_eng=round(sum(v["engagement_rate"] for v in d["videos"])/max(len(d["videos"]),1),2)
        vals=[f"{medals[ri]} #{ri+1}",name,str(score),_fmt(d["subscriber_count"]),
              _fmt(d["avg_views"]),str(d["upload_freq_per_week"]),f"{avg_eng}%"]
        for val,x,w in zip(vals,x_pos,col_w):
            _add_textbox(sl,val,x,y,w-0.1,0.6,11,ri==0,WHITE if ri>0 else ACCENT,PP_ALIGN.CENTER)

# ── Main builder ─────────────────────────────────────────────────────────────

def generate_pptx(all_data, your_company) -> bytes:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    scores = _rank_companies(all_data)
    slide_cover(prs, all_data, your_company)
    slide_executive_summary(prs, all_data, scores)
    slide_channel_overview(prs, all_data)
    slide_content_performance(prs, all_data)
    slide_topics_themes(prs, all_data)
    slide_posting_frequency(prs, all_data)
    slide_engagement(prs, all_data)
    slide_gap_analysis(prs, all_data)
    slide_recommendations(prs, all_data, your_company)
    slide_summary_ranking(prs, all_data, scores)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
